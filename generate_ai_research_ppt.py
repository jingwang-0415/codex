from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTFILE = Path("/Users/a747/codex/全球AI提升软件工程开发效率调研报告_2026-04-20.pptx")


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


BG = RGBColor(248, 244, 236)
NAVY = RGBColor(21, 45, 78)
TEAL = RGBColor(39, 110, 120)
ORANGE = RGBColor(202, 113, 44)
INK = RGBColor(45, 52, 54)
MUTED = RGBColor(96, 108, 118)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(235, 240, 242)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.8), Inches(0.9))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = NAVY
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.72), Inches(1.1), Inches(11.5), Inches(0.45))
        tf2 = box2.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(10)
        r2.font.color.rgb = MUTED


def add_footer(slide, text="资料来源均来自公开官方报告、官方博客或官方客户案例"):
    box = slide.shapes.add_textbox(Inches(0.75), Inches(7.0), Inches(12.0), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(8)
    r.font.color.rgb = MUTED


def add_accent(slide, color=TEAL):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(0.35), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bullets(slide, items, left=0.95, top=1.7, width=11.2, height=4.7, font_size=20, color=INK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
        p.line_spacing = 1.15


def add_two_col(slide, left_title, left_items, right_title, right_items):
    for x, title, items, color in [
        (0.8, left_title, left_items, TEAL),
        (6.75, right_title, right_items, ORANGE),
    ]:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.65), Inches(5.55), Inches(4.85)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = LIGHT

        title_box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(1.88), Inches(4.9), Inches(0.4))
        p = title_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = color

        add_bullets(
            slide,
            items,
            left=x + 0.22,
            top=2.35,
            width=4.95,
            height=3.9,
            font_size=15,
            color=INK,
        )


def add_case_grid(slide, cases):
    positions = [
        (0.8, 1.65), (4.45, 1.65), (8.1, 1.65),
        (0.8, 4.2), (4.45, 4.2), (8.1, 4.2),
    ]
    for (x, y), case in zip(positions, cases):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.15), Inches(1.95)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT

        title_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.14), Inches(2.7), Inches(0.35))
        p = title_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = case["name"]
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = NAVY

        tool_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.5), Inches(2.8), Inches(0.3))
        p2 = tool_box.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = f"工具：{case['tool']}"
        r2.font.size = Pt(10)
        r2.font.color.rgb = TEAL

        body = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.82), Inches(2.75), Inches(0.95))
        tf = body.text_frame
        tf.word_wrap = True
        for i, line in enumerate(case["points"]):
            p3 = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p3.text = line
            p3.font.size = Pt(10)
            p3.font.color.rgb = INK
            p3.space_after = Pt(3)


def add_banner(slide, text, color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.65), Inches(11.7), Inches(0.55)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    box = slide.shapes.add_textbox(Inches(1.0), Inches(1.76), Inches(11.0), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = WHITE


def add_source_box(slide, lines):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(5.75), Inches(11.3), Inches(0.85)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LIGHT
    box = slide.shapes.add_textbox(Inches(1.15), Inches(5.95), Inches(10.8), Inches(0.45))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(9)
        p.font.color.rgb = MUTED
        p.space_after = Pt(2)


def add_case_detail(slide, name, tool, before, process, after, proof, howto, source):
    add_banner(slide, f"{name} | 工具：{tool}", NAVY)
    cards = [
        ("使用前", before, ORANGE, 0.8, 2.35, 3.0, 1.7),
        ("使用过程", process, TEAL, 4.15, 2.35, 4.9, 1.7),
        ("使用后对比", after, NAVY, 9.35, 2.35, 3.1, 1.7),
        ("证明了什么", proof, TEAL, 0.8, 4.35, 5.8, 1.75),
        ("怎么借鉴", howto, ORANGE, 6.9, 4.35, 5.55, 1.75),
    ]
    for title, items, color, x, y, w, h in cards:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = LIGHT
        tbox = slide.shapes.add_textbox(Inches(x + 0.16), Inches(y + 0.12), Inches(w - 0.3), Inches(0.28))
        p = tbox.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = color
        add_bullets(slide, items, left=x + 0.16, top=y + 0.43, width=w - 0.3, height=h - 0.5, font_size=10, color=INK)
    add_source_box(slide, [source])


def new_slide(title, subtitle=None, accent=TEAL):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_accent(slide, accent)
    add_title(slide, title, subtitle)
    add_footer(slide)
    return slide


# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
hero = slide.shapes.add_shape(
    MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
)
hero.fill.solid()
hero.fill.fore_color.rgb = NAVY
hero.line.fill.background()

panel = slide.shapes.add_shape(
    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(0.95), Inches(11.85), Inches(5.35)
)
panel.fill.solid()
panel.fill.fore_color.rgb = BG
panel.line.fill.background()

box = slide.shapes.add_textbox(Inches(1.2), Inches(1.55), Inches(10.8), Inches(1.8))
tf = box.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
r = p.add_run()
r.text = "全球 AI 提升软件工程开发效率\n现状与案例调研"
r.font.size = Pt(28)
r.font.bold = True
r.font.color.rgb = NAVY

sub = slide.shapes.add_textbox(Inches(1.25), Inches(3.55), Inches(9.8), Inches(0.7))
p2 = sub.text_frame.paragraphs[0]
r2 = p2.add_run()
r2.text = "基于公开官方报告、客户案例与开发者调查整理"
r2.font.size = Pt(16)
r2.font.color.rgb = MUTED

meta = slide.shapes.add_textbox(Inches(1.25), Inches(5.55), Inches(5.5), Inches(0.4))
p3 = meta.text_frame.paragraphs[0]
r3 = p3.add_run()
r3.text = "日期：2026-04-20"
r3.font.size = Pt(12)
r3.font.color.rgb = ORANGE


# Slide 2
slide = new_slide("核心结论", "AI 已从代码补全演进为覆盖整个软件工程流程的协作式与代理式工作流", TEAL)
add_bullets(
    slide,
    [
        "使用率高：AI 已成为开发团队主流工具，而非可选项。",
        "重点场景成熟：代码生成、测试、调试、文档、评审、迁移升级最先出 ROI。",
        "组织差异决定效果：统一平台、测试体系、反馈回路和治理能力会放大收益。",
        "趋势升级明显：从 IDE 内 Copilot，走向接入代码库、工单、日志和 CI 的工程代理。",
        "管理重点改变：效率提升不再只看写码速度，而看吞吐、质量、稳定性是否一起改善。",
    ],
)


# Slide 3
slide = new_slide("现状 1：使用率已高，但信任仍不足", "对应调研文档 2.1", TEAL)
add_banner(slide, "主流组织已经在用 AI，但开发者仍然默认需要人工验证", TEAL)
add_two_col(
    slide,
    "公开数据",
    [
        "84% 的开发者已使用或计划使用 AI 工具参与开发。",
        "50.6% 的职业开发者表示每天使用 AI。",
        "46% 的开发者不信任 AI 输出准确性。",
        "只有 33% 的开发者表示信任 AI 输出。",
    ],
    "管理含义",
    [
        "AI 已经普及，但尚未获得“默认可信”的工程地位。",
        "高频使用不代表可跳过测试、评审和验收。",
        "团队要把验证体系建设成 AI 提效的前提，而不是事后补救。",
        "最有效的定位是“加速器 + 建议者”，不是完全替代者。",
    ],
)
add_source_box(
    slide,
    [
        "来源：Stack Overflow Developer Survey 2025 AI Section",
        "链接：https://survey.stackoverflow.co/2025/ai",
    ],
)


# Slide 4
slide = new_slide("现状 2：最成熟的场景集中在高频、重复、可验证任务", "对应调研文档 2.2", ORANGE)
add_banner(slide, "AI 最容易先替代机械性劳动，再辅助复杂工程决策", ORANGE)
add_two_col(
    slide,
    "成熟场景",
    [
        "代码生成与补全：函数、样板代码、接口、脚手架。",
        "代码理解：解释陌生仓库、定位依赖、生成说明。",
        "测试与调试：生成单测、定位报错、给出修复建议。",
        "文档沉淀：PR 摘要、变更说明、技术文档。",
    ],
    "为什么这些场景先成熟",
    [
        "任务重复率高，标准化程度高。",
        "结果可被编译、测试、lint、规则快速验证。",
        "依赖的上下文相对明确，易于接入代码库与文档。",
        "失败成本相对可控，适合逐步扩大权限。",
    ],
)


# Slide 5
slide = new_slide("现状 3：高责任环节仍然谨慎", "对应调研文档 2.3", NAVY)
add_banner(slide, "在部署、监控、项目规划等高责任场景，全球团队仍更保守", NAVY)
add_two_col(
    slide,
    "公开调查",
    [
        "75.8% 的受访者不计划把 AI 用于 deployment and monitoring。",
        "69.2% 的受访者不计划把 AI 用于 project planning。",
        "说明 AI 在高风险、高责任环节仍未获得充分授权。",
    ],
    "实际启示",
    [
        "AI 更适合作为建议者、自动化执行者和辅助分析者。",
        "最终上线、资源权衡、优先级判断仍需人工拍板。",
        "若要扩展到高责任环节，必须同步增加权限治理和回滚机制。",
    ],
)
add_source_box(
    slide,
    [
        "来源：Stack Overflow Developer Survey 2025 AI Section",
        "链接：https://survey.stackoverflow.co/2025/ai",
    ],
)


# Slide 6
slide = new_slide("现状 4：AI 更像组织能力放大器", "对应调研文档 2.4", TEAL)
add_banner(slide, "能否真正提效，关键不在“有没有 AI”，而在“工程体系是否接得住”", TEAL)
add_two_col(
    slide,
    "DORA 2025 观察",
    [
        "90% 的技术从业者表示已在工作中使用 AI。",
        "超过 80% 认为 AI 提高了生产力。",
        "AI 与吞吐量、产品表现正相关。",
        "AI 与交付稳定性仍存在负相关风险。",
    ],
    "组织要求",
    [
        "需要高质量内部平台、清晰工作流和快速反馈环路。",
        "自动化测试、版本管理和可观测性是必要护栏。",
        "如果质量体系薄弱，AI 往往只会放大现有问题。",
        "先进团队更重视系统效率，而非单点写码速度。",
    ],
)
add_source_box(
    slide,
    [
        "来源：DORA 2025 State of AI-assisted Software Development / Google Cloud",
        "链接：https://dora.dev/report/2025",
    ],
)


# Slide 7
slide = new_slide("当前主流落地模式", "从工具辅助走向代理协作", ORANGE)
add_bullets(
    slide,
    [
        "IDE 内 Copilot：在 VS Code / JetBrains 中做补全、解释、重构、问答。",
        "代码库理解与跨文件修改：对整个仓库建立索引，处理依赖更新和多文件修改。",
        "代理式开发：读取代码、文档、Issue、日志，自动改代码、跑测试、开 PR。",
        "AI 代码评审：在 PR 阶段自动检查逻辑缺陷、边界场景和安全问题。",
        "现代化改造：面向 Java/.NET 升级、框架迁移、技术债治理等高 ROI 场景。",
        "治理驱动模式：围绕接受率、PR 周转、缺陷率和稳定性建立度量与护栏。",
    ],
    font_size=18,
)


# Slide 8
slide = new_slide("代表案例一览", "文档中已展开 7 个案例，这里汇总最具代表性的 6 个", TEAL)
cases = [
    {"name": "Accenture", "tool": "GitHub Copilot", "points": ["企业大规模发放许可证", "日常编码中接受建议", "按接受率与合并率评估效果"]},
    {"name": "Carlsberg", "tool": "Copilot + GHAS", "points": ["统一研发平台", "把 AI 接进安全与 CI 流程", "减少工具切换成本"]},
    {"name": "Dun & Bradstreet", "tool": "Gemini Code Assist", "points": ["先评估再定制接入", "用于代码建议与评审", "兼顾隐私与治理"]},
    {"name": "Novacomp / Signaturit", "tool": "Amazon Q Transform", "points": ["以代理方式做版本迁移", "自动更新代码并验证", "现代化改造 ROI 明确"]},
    {"name": "NVIDIA", "tool": "Cursor", "points": ["把 AI 扩展到整个 SDLC", "接入文档、票据与测试", "形成代理式工作流"]},
    {"name": "PlanetScale", "tool": "Cursor Bugbot", "points": ["聚焦 PR 阶段质量控制", "自动识别高价值缺陷", "缓解评审瓶颈"]},
]
add_case_grid(slide, cases)


# Slide 9
slide = new_slide("汇报建议", "如果用于内部汇报，可重点落在这 4 个判断", NAVY)
add_bullets(
    slide,
    [
        "第一，全球已经从“AI 辅助写代码”走向“AI 参与整个软件工程流程”。",
        "第二，领先组织的优势不只在模型，而在上下文接入、自动验证和治理闭环。",
        "第三，最容易先出 ROI 的不是最前沿代理，而是高频、重复、可验证任务。",
        "第四，衡量 AI 成效的标准正在从“用了多少”转向“是否改善吞吐、质量和稳定性”。",
    ],
    top=2.0,
    height=3.6,
    font_size=21,
)
add_source_box(
    slide,
    [
        "建议与文档配套使用：详细出处、工具、AI 使用过程、效果指标已写入研究文档。",
        "文档路径：/Users/a747/codex/全球AI提升软件工程开发效率调研报告_2026-04-20.md",
    ],
)


# Slide 10
slide = new_slide("详细案例 1", "Accenture：企业如何验证 Copilot 的真实收益", TEAL)
add_case_detail(
    slide,
    "Accenture",
    "GitHub Copilot",
    [
        "企业需要先证明 AI 真能提升交付效率。",
        "推广前缺少可量化的企业级证据。",
    ],
    [
        "先用 20 人试点，再做随机对照试验。",
        "450 名使用者对比 200 名未使用者。",
        "发放许可证并安装 IDE 扩展。",
        "跟踪接受率、提交、PR 与构建指标。",
    ],
    [
        "67% 每周至少使用 5 天。",
        "30% 建议被接受。",
        "90% 提交过含 AI 代码的提交。",
        "91% 团队合并过含 AI 代码的 PR。",
    ],
    [
        "Copilot 最适合作为现有开发流上的增效层。",
        "企业完全可以用工程指标而非主观感受评估 AI。",
    ],
    [
        "先做小规模试点，不直接全量推广。",
        "指标优先看安装率、活跃率、接受率、PR 周转。",
    ],
    "来源：GitHub 官方研究与 GitHub Customer Story（Accenture）",
)


# Slide 11
slide = new_slide("详细案例 2", "Carlsberg：先统一平台，再让 AI 真正起效", ORANGE)
add_case_detail(
    slide,
    "Carlsberg",
    "Copilot + GHAS + Actions",
    [
        "研发流程分散在 6 套工具中。",
        "上下文切换频繁，源码中有未识别 secrets。",
    ],
    [
        "先统一到 GitHub Enterprise。",
        "再接入 Copilot、GHAS、Secret Scanning 和 Actions。",
        "让生成代码继续经过扫描、CI 和合并控制。",
    ],
    [
        "6 套工具收敛到 1 个统一平台。",
        "发现约 600 个未知 secrets。",
        "移除 30,000+ 安全漏洞。",
    ],
    [
        "AI 收益往往依赖统一平台与安全护栏。",
        "减少上下文切换本身就是效率收益来源。",
    ],
    [
        "工具分散的团队先做平台整合。",
        "AI 工具和安全、CI 一起上，而不是单独上。",
    ],
    "来源：GitHub Customer Story（Carlsberg Group）",
)


# Slide 12
slide = new_slide("详细案例 3", "Dun & Bradstreet：在高安全要求下落地 AI 开发", TEAL)
add_case_detail(
    slide,
    "Dun & Bradstreet",
    "Gemini Code Assist",
    [
        "onboarding 慢，遗留知识孤岛明显。",
        "测试覆盖低，现代化改造困难。",
    ],
    [
        "把 AI 引入定义为 people/process/tools 联动改造。",
        "评估开源与商业方案后选 Gemini。",
        "与 Google Cloud 共同配置并接入原有 review 流。",
        "用于代码建议、review、lint、Kotlin 转换和新人辅导。",
    ],
    [
        "早期内部指标显示生产率提升约 30%。",
        "Spring 到 Kotlin 转换节省时间。",
        "问题更早发现，新人 ramp up 更快。",
    ],
    [
        "AI 不只用于写代码，也能提升 onboarding 和知识共享。",
        "高安全行业更适合先设计护栏再扩大推广。",
    ],
    [
        "试点优先选 review、代码解释、现代化和 onboarding。",
        "把隐私、日志、管理控制与安全扫描一起纳入方案。",
    ],
    "来源：Google Cloud 官方客户案例（Dun & Bradstreet）",
)


# Slide 13
slide = new_slide("详细案例 4", "Amazon Q：把现代化改造做成可量化 ROI 项目", NAVY)
add_case_detail(
    slide,
    "Novacomp / Signaturit",
    "Amazon Q Developer Transform",
    [
        "Java/.NET 遗留系统升级长期被挤压。",
        "这类改造高度依赖高级工程师，人力昂贵。",
    ],
    [
        "在 IDE 或 CLI 发起 transformation 任务。",
        "代理生成计划并执行升级与迁移。",
        "输出 diff、summary、测试和部署准备结果。",
        "人工审查后接受并合并。",
    ],
    [
        "Novacomp：3 周缩短到 50 分钟。",
        "平均技术债下降 60%。",
        "Signaturit：6-8 个月缩短为几天。",
    ],
    [
        "版本升级和迁移是当前最容易算清 ROI 的 AI 场景。",
        "代理式改造必须配套计划、diff 和验证结果。",
    ],
    [
        "首批项目选中等规模且容易验证的系统。",
        "固定要求：计划、diff、测试、人工验收。",
    ],
    "来源：AWS 官方案例、产品页与文档",
)


# Slide 14
slide = new_slide("详细案例 5", "NVIDIA：从写代码走向 SDLC 自动化", TEAL)
add_case_detail(
    slide,
    "NVIDIA",
    "Cursor",
    [
        "代码库巨大且跨栈依赖复杂。",
        "此前其他 AI 工具未带来显著提升。",
    ],
    [
        "先用 Cursor 做大型代码库语义理解。",
        "再扩展到测试、评审、调试和工作流自动化。",
        "通过自定义规则自动化 git flow。",
        "用 MCP 接 tickets 与文档，再修复并跑测试。",
    ],
    [
        "30,000+ 开发者每天使用。",
        "提交代码量提升到 3 倍以上。",
        "新人 onboarding 更快。",
    ],
    [
        "复杂工程的关键不只是模型，而是上下文接入能力。",
        "AI 的成熟形态是自动执行工作流，而不是只补全代码。",
    ],
    [
        "优先接代码库、工单、文档和 CI。",
        "先从 bug 修复链路做半自动闭环。",
    ],
    "来源：Cursor 官方客户案例（NVIDIA）",
)


# Slide 15
slide = new_slide("详细案例 6", "PlanetScale：AI 写码更快后，评审成为新瓶颈", ORANGE)
add_case_detail(
    slide,
    "PlanetScale",
    "Cursor Bugbot",
    [
        "代码生成变快，但人工 review 能力不变。",
        "如果只靠人力扩容，需要额外 2 名工程师专做评审。",
    ],
    [
        "在 PR 流程引入 Bugbot 作为 agentic review layer。",
        "重点检查状态同步、逻辑流、异步交互和边界条件。",
        "以 merge time 解决率衡量评论价值。",
    ],
    [
        "每月评审 2,000+ PR。",
        "约 80% 评论在合并前被处理。",
        "节省约 2 FTE 的 review 工作量。",
    ],
    [
        "AI 提高吞吐后，下游 review 会成为瓶颈。",
        "专门的 AI review 层比直接用通用模型审 PR 更有效。",
    ],
    [
        "如果团队已大量用 AI 写代码，下一步就该补 AI review。",
        "指标看高价值评论解决率，而不是评论数量。",
    ],
    "来源：Cursor 官方客户案例（PlanetScale）",
)


# Slide 16
slide = new_slide("详细案例 7", "OpenAI：agent-first 软件工程如何工作", NAVY)
add_case_detail(
    slide,
    "OpenAI",
    "Codex",
    [
        "从空仓库出发，设定 0 行人工代码的极端约束。",
        "核心问题转为如何设计环境与反馈回路。",
    ],
    [
        "工程师用 prompt 定义目标与约束。",
        "Codex 生成应用、测试、CI、文档和可观测性。",
        "仓库知识成为 system of record。",
        "把日志、指标、浏览器工具暴露给代理做验证。",
    ],
    [
        "5 个月交付内部 beta 产品。",
        "约 100 万行代码、1,500 个 PR。",
        "整体速度约为人工方式的 1/10。",
    ],
    [
        "agent-first 的关键是环境、约束和反馈回路，而不只是模型。",
        "工程师角色从写代码转向定义目标与治理代理。",
    ],
    [
        "大多数团队不应直接一步到位。",
        "先补测试、CI、仓库文档，再尝试代理闭环。",
    ],
    "来源：OpenAI 官方工程文章（Harness engineering）",
)


# Slide 17
slide = new_slide("结论与证明", "把结论和案例证据放在一起，便于直接汇报", TEAL)
add_two_col(
    slide,
    "关键结论",
    [
        "高频、重复、可验证任务最先出 ROI。",
        "上下文接入能力比模型本身更关键。",
        "编码提速后，review 与验证会成为新瓶颈。",
        "现代化改造最容易算清 ROI。",
    ],
    "对应案例",
    [
        "Accenture、D&B、Amazon Q",
        "NVIDIA、OpenAI、Carlsberg",
        "PlanetScale、DORA 2025",
        "Novacomp、Signaturit",
    ],
)
add_source_box(
    slide,
    [
        "可搭配文档第 5 节《结论与证据链》使用。",
        "文档路径：/Users/a747/codex/全球AI提升软件工程开发效率调研报告_2026-04-20.md",
    ],
)


# Slide 18
slide = new_slide("落地建议", "更适合大多数组织的 4 阶段导入路径", ORANGE)
add_bullets(
    slide,
    [
        "第一阶段 Copilot 化：先做代码补全、单测、文档和代码解释。",
        "第二阶段 上下文化：接入代码库、需求票据、文档、CI 和测试结果。",
        "第三阶段 验证化：补 AI review、自动测试、安全扫描和 merge gate。",
        "第四阶段 代理化：让 AI 处理更长链路任务，人工转向目标定义与验收。",
        "如果要选首批试点，优先：代码解释、单测、版本升级、遗留系统现代化。",
    ],
    top=1.95,
    height=4.2,
    font_size=18,
)
add_source_box(
    slide,
    [
        "可搭配文档第 6 节《如何落地使用 AI 提升软件工程效率》使用。",
        "建议汇报时把案例、证据和实施路径一起呈现。",
    ],
)


prs.save(str(OUTFILE))
print(f"Generated: {OUTFILE}")
